#!/usr/bin/env python3
"""
JJクラブ 大曽根 - スマホLP デザイン PPTX生成スクリプト
参考: バニラ求人ページ (https://tokai.qzin.jp/showtikubai/)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== カラーパレット =====
PINK_MAIN = RGBColor(0xE9, 0x1E, 0x63)      # メインピンク（CTA・アクセント）
PINK_LIGHT = RGBColor(0xFC, 0xE4, 0xEC)      # 薄ピンク（背景）
PINK_DARK = RGBColor(0xC2, 0x18, 0x5B)       # 濃ピンク
GOLD = RGBColor(0xD4, 0xA0, 0x17)            # ゴールド（高級感）
GOLD_LIGHT = RGBColor(0xFD, 0xF0, 0xD5)      # 薄ゴールド
BLACK = RGBColor(0x1A, 0x1A, 0x2E)           # 濃紺黒
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)       # 濃グレー
MID_GRAY = RGBColor(0x66, 0x66, 0x66)        # 中間グレー
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)      # 薄グレー背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_LINE = RGBColor(0x06, 0xC7, 0x55)      # LINE緑
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)          # 紫アクセント

# ===== スマホサイズ設定 (9:16比率) =====
SLIDE_WIDTH = Cm(10.8)   # 約4.25インチ
SLIDE_HEIGHT = Cm(19.2)  # 約7.56インチ

# マージン
MARGIN_X = Cm(0.8)
CONTENT_WIDTH = Cm(9.2)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    # 空のレイアウトを使用
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs, blank_layout


def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    """背景色の矩形を追加"""
    w = width or SLIDE_WIDTH
    h = height or SLIDE_HEIGHT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=12, font_color=DARK_GRAY, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name='Yu Gothic UI',
                 line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=12, font_color=DARK_GRAY, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name='Yu Gothic UI',
                       line_spacing=1.4):
    """複数行テキストボックスを追加（各行は (text, size, color, bold) のタプル可）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, sz, col, b = line
        else:
            text, sz, col, b = line, font_size, font_color, bold

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = col
        p.font.bold = b
        p.font.name = font_name
        p.alignment = alignment
        p.line_spacing = Pt(sz * line_spacing)
        p.space_after = Pt(2)

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=True, radius=Cm(0.3)):
    """角丸ボタン/ラベルを追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # 角丸の調整
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.3  # 角丸の度合い

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Yu Gothic UI'
        p.alignment = PP_ALIGN.CENTER
        tf.auto_size = None
    return shape


def add_image_placeholder(slide, left, top, width, height, label="画像"):
    """画像プレースホルダー（グレー矩形＋テキスト）"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p.font.bold = False
    p.font.name = 'Yu Gothic UI'
    p.alignment = PP_ALIGN.CENTER

    return shape


def add_divider(slide, top, color=RGBColor(0xEE, 0xEE, 0xEE)):
    """水平区切り線"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(1.5), top, Cm(7.8), Pt(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle_number(slide, left, top, number, color=PINK_MAIN):
    """番号付き丸を追加"""
    size = Cm(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Yu Gothic UI'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ============================================================
# スライド1: ファーストビュー（ヒーロー）
# ============================================================
def slide_hero(prs, layout):
    slide = prs.slides.add_slide(layout)

    # 背景画像プレースホルダー（フルスクリーン）
    add_image_placeholder(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
                          "キービジュアル画像\n（店舗イメージ写真）\n1080 x 1920px")

    # 半透明オーバーレイ（下部グラデーション風）
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(10), SLIDE_WIDTH, Cm(9.2))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    overlay.line.fill.background()
    # 透明度設定（XMLレベル）
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = overlay._element.spPr
    solidFill_el = spPr.find(qn('a:solidFill'))
    if solidFill_el is None:
        # solidFillがない場合はfillの中を探す
        for el in spPr.iter():
            if el.tag.endswith('srgbClr'):
                alpha = etree.SubElement(el, qn('a:alpha'))
                alpha.set('val', '60000')
                break

    # ロゴ/店名ラベル
    add_rounded_rect(slide, Cm(1.5), Cm(1.0), Cm(7.8), Cm(1.2),
                     PINK_MAIN, "JJクラブ 大曽根", font_size=16, bold=True)

    # メインキャッチ
    add_multiline_text(slide, [
        ("業界未経験でも", 22, WHITE, True),
        ("他の店舗からの移籍でも", 22, WHITE, True),
        ("あなたをしっかり", 22, WHITE, True),
        ("サポートします！", 26, RGBColor(0xFF, 0xD5, 0x4F), True),
    ], Cm(0.8), Cm(11.0), Cm(9.2), Cm(5.5), alignment=PP_ALIGN.CENTER)

    # サブキャッチ
    add_multiline_text(slide, [
        ("身バレ完全回避のお店！顔出しナシ！", 13, WHITE, False),
        ("繁華街じゃないから目立たず出勤できる！", 13, WHITE, False),
        ("なのに総合駅近くだから集客力抜群！", 13, WHITE, False),
    ], Cm(0.8), Cm(15.0), Cm(9.2), Cm(2.5), alignment=PP_ALIGN.CENTER)

    # CTAボタン
    add_rounded_rect(slide, Cm(1.5), Cm(17.2), Cm(7.8), Cm(1.3),
                     PINK_MAIN, "今すぐ応募する ▶", font_size=16, bold=True)

    return slide


# ============================================================
# スライド2: 当店で働く3つのメリット
# ============================================================
def slide_merits(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, WHITE)

    # セクションヘッダー
    header_bg = add_bg_rect(slide, PINK_MAIN, 0, 0, SLIDE_WIDTH, Cm(2.2))
    add_text_box(slide, "当店で働く 3つのメリット", Cm(0.8), Cm(0.5), Cm(9.2), Cm(1.2),
                 font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    y_start = Cm(2.8)

    # メリット1
    add_circle_number(slide, Cm(0.8), y_start, "1", PINK_MAIN)
    add_text_box(slide, "未経験者大歓迎！", Cm(2.3), y_start, Cm(7.5), Cm(0.8),
                 font_size=16, font_color=PINK_DARK, bold=True)
    add_text_box(slide,
                 "初めてでも安心のサポート体制。マンツーマンで丁寧に教えます。講習制度もあるので、何も分からなくても大丈夫です！",
                 Cm(0.8), y_start + Cm(1.3), Cm(9.2), Cm(2.0),
                 font_size=11, font_color=MID_GRAY)

    # 画像スペース
    add_image_placeholder(slide, Cm(1.5), y_start + Cm(3.2), Cm(7.8), Cm(2.5),
                          "サポートイメージ画像")

    y2 = y_start + Cm(6.2)
    add_divider(slide, y2 - Cm(0.3))

    # メリット2
    add_circle_number(slide, Cm(0.8), y2, "2", PINK_MAIN)
    add_text_box(slide, "完全自由出勤制！", Cm(2.3), y2, Cm(7.5), Cm(0.8),
                 font_size=16, font_color=PINK_DARK, bold=True)
    add_text_box(slide,
                 "出勤日・時間は全てあなた次第。週1日〜OK！掛け持ちもOK！あなたのライフスタイルに合わせて自由に働けます。",
                 Cm(0.8), y2 + Cm(1.3), Cm(9.2), Cm(2.0),
                 font_size=11, font_color=MID_GRAY)

    y3 = y2 + Cm(3.5)
    add_divider(slide, y3 - Cm(0.3))

    # メリット3
    add_circle_number(slide, Cm(0.8), y3, "3", PINK_MAIN)
    add_text_box(slide, "充実の待遇！", Cm(2.3), y3, Cm(7.5), Cm(0.8),
                 font_size=16, font_color=PINK_DARK, bold=True)
    add_multiline_text(slide, [
        ("✓ アリバイ対策万全", 11, MID_GRAY, False),
        ("✓ 完全個室待機", 11, MID_GRAY, False),
        ("✓ 送迎あり", 11, MID_GRAY, False),
        ("✓ 日払い対応", 11, MID_GRAY, False),
        ("✓ 各種手当充実", 11, MID_GRAY, False),
    ], Cm(0.8), y3 + Cm(1.3), Cm(9.2), Cm(3.5), font_size=11)

    return slide


# ============================================================
# スライド3: 体験入店情報
# ============================================================
def slide_trial(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, GOLD_LIGHT)

    # ヘッダー
    add_rounded_rect(slide, Cm(1.5), Cm(0.8), Cm(7.8), Cm(1.5),
                     GOLD, "体験入店", font_size=20, bold=True)

    # メイン金額
    add_text_box(slide, "本日の体験入店", Cm(0.8), Cm(3.0), Cm(9.2), Cm(0.8),
                 font_size=14, font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    add_multiline_text(slide, [
        ("日給保証", 14, DARK_GRAY, False),
        ("60,000円", 36, PINK_MAIN, True),
    ], Cm(0.8), Cm(3.8), Cm(9.2), Cm(3.5),
       font_size=36, alignment=PP_ALIGN.CENTER)

    # 募集人数
    add_rounded_rect(slide, Cm(2.5), Cm(6.2), Cm(5.8), Cm(1.0),
                     PINK_MAIN, "本日の募集人数：5人", font_size=13, bold=True)

    # 時間別給与
    add_text_box(slide, "時間別 給与例", Cm(0.8), Cm(7.8), Cm(9.2), Cm(0.7),
                 font_size=14, font_color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    y_table = Cm(8.8)
    time_data = [
        ("4時間", "40,000円"),
        ("7時間", "60,000円"),
        ("10時間", "80,000円"),
    ]

    for i, (time_str, amount) in enumerate(time_data):
        row_y = y_table + Cm(i * 1.6)
        # 時間ラベル
        add_rounded_rect(slide, Cm(1.0), row_y, Cm(3.5), Cm(1.2),
                         WHITE, time_str, font_size=13, font_color=DARK_GRAY, bold=True)
        # 金額
        add_text_box(slide, amount, Cm(5.0), row_y, Cm(4.8), Cm(1.2),
                     font_size=20, font_color=PINK_MAIN, bold=True, alignment=PP_ALIGN.CENTER)

    # キービジュアルスペース
    add_image_placeholder(slide, Cm(1.5), Cm(14.2), Cm(7.8), Cm(3.0),
                          "体験入店イメージ画像")

    # CTA
    add_rounded_rect(slide, Cm(1.5), Cm(17.6), Cm(7.8), Cm(1.2),
                     PINK_MAIN, "体験入店に応募する ▶", font_size=14, bold=True)

    return slide


# ============================================================
# スライド4: 安心ポイント（身バレ対策）
# ============================================================
def slide_safety(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, WHITE)

    # ヘッダー
    header_bg = add_bg_rect(slide, BLACK, 0, 0, SLIDE_WIDTH, Cm(2.5))
    add_multiline_text(slide, [
        ("身バレ完全回避！", 20, WHITE, True),
        ("安心のプライバシー対策", 14, RGBColor(0xFF, 0xD5, 0x4F), True),
    ], Cm(0.8), Cm(0.4), Cm(9.2), Cm(2.0), alignment=PP_ALIGN.CENTER)

    y = Cm(3.2)

    # 安心ポイントリスト
    points = [
        ("顔出しナシ！", "HP写真はモザイク加工。身バレの心配はありません。"),
        ("繁華街じゃないから安心", "大曽根エリアだから、目立たず出勤できます。"),
        ("完全個室待機", "他の女の子と顔を合わせることなく待機できます。"),
        ("アリバイ対策完備", "給料明細の発行社名変更、架電対応など万全のサポート。"),
        ("マジックミラー完備", "お客様をミラー越しに確認。NGなら断れます。"),
        ("女性スタッフ在籍", "女性スタッフが受付対応。安心してご相談ください。"),
    ]

    for i, (title, desc) in enumerate(points):
        row_y = y + Cm(i * 2.3)

        # アイコン的な丸
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(0.8), row_y + Cm(0.1), Cm(0.8), Cm(0.8))
        icon.fill.solid()
        icon.fill.fore_color.rgb = PINK_LIGHT
        icon.line.fill.background()
        tf = icon.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(10)
        p.font.color.rgb = PINK_MAIN
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, title, Cm(2.0), row_y, Cm(8.0), Cm(0.7),
                     font_size=13, font_color=DARK_GRAY, bold=True)
        add_text_box(slide, desc, Cm(2.0), row_y + Cm(0.7), Cm(8.0), Cm(1.2),
                     font_size=10, font_color=MID_GRAY)

    return slide


# ============================================================
# スライド5: 給与情報
# ============================================================
def slide_salary(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, PINK_LIGHT)

    # ヘッダー
    add_rounded_rect(slide, Cm(1.5), Cm(0.8), Cm(7.8), Cm(1.5),
                     PINK_MAIN, "給与について", font_size=20, bold=True)

    # メイン給与情報
    add_text_box(slide, "60分コース", Cm(0.8), Cm(3.0), Cm(9.2), Cm(0.8),
                 font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    add_multiline_text(slide, [
        ("バック", 14, DARK_GRAY, False),
        ("13,000円以上", 32, PINK_MAIN, True),
    ], Cm(0.8), Cm(3.8), Cm(9.2), Cm(3.0), alignment=PP_ALIGN.CENTER)

    # 補足
    add_text_box(slide, "指名本数スライド制でさらにアップ！",
                 Cm(0.8), Cm(6.5), Cm(9.2), Cm(0.8),
                 font_size=12, font_color=PINK_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # 区切り線
    add_divider(slide, Cm(7.5), PINK_MAIN)

    # 収入例
    add_text_box(slide, "実際の収入例", Cm(0.8), Cm(8.0), Cm(9.2), Cm(0.7),
                 font_size=15, font_color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # 例1
    box1 = add_bg_rect(slide, WHITE, Cm(0.8), Cm(9.0), Cm(9.2), Cm(3.0))
    box1.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    box1.line.width = Pt(1)
    add_multiline_text(slide, [
        ("Aさん（30歳・パート主婦）", 12, PINK_DARK, True),
        ("7時間勤務 / 4本接客", 11, MID_GRAY, False),
        ("日収：51,000円", 18, PINK_MAIN, True),
        ("月収例（週3日）：約60万円", 12, DARK_GRAY, True),
    ], Cm(1.2), Cm(9.2), Cm(8.5), Cm(2.8), alignment=PP_ALIGN.CENTER)

    # 例2
    box2 = add_bg_rect(slide, WHITE, Cm(0.8), Cm(12.5), Cm(9.2), Cm(3.0))
    box2.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    box2.line.width = Pt(1)
    add_multiline_text(slide, [
        ("Bさん（25歳・OL掛け持ち）", 12, PINK_DARK, True),
        ("5時間勤務 / 3本接客", 11, MID_GRAY, False),
        ("日収：39,000円", 18, PINK_MAIN, True),
        ("月収例（週2日）：約30万円", 12, DARK_GRAY, True),
    ], Cm(1.2), Cm(12.7), Cm(8.5), Cm(2.8), alignment=PP_ALIGN.CENTER)

    # 注記
    add_text_box(slide, "※上記は一例です。接客数・コースにより変動します",
                 Cm(0.8), Cm(16.0), Cm(9.2), Cm(0.7),
                 font_size=9, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # CTA
    add_rounded_rect(slide, Cm(1.5), Cm(17.0), Cm(7.8), Cm(1.2),
                     PINK_MAIN, "給与について詳しく聞く ▶", font_size=13, bold=True)

    return slide


# ============================================================
# スライド6: 在籍女性の声（インタビュー）
# ============================================================
def slide_interview(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, WHITE)

    # ヘッダー
    header_bg = add_bg_rect(slide, PURPLE, 0, 0, SLIDE_WIDTH, Cm(2.2))
    add_text_box(slide, "在籍女性の声", Cm(0.8), Cm(0.5), Cm(9.2), Cm(1.2),
                 font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # インタビュー1
    y1 = Cm(2.8)
    add_image_placeholder(slide, Cm(0.8), y1, Cm(3.0), Cm(3.5),
                          "女性写真1\n（モザイク）")

    add_multiline_text(slide, [
        ("Cさん（22歳・未経験）", 12, PURPLE, True),
        ("入店3ヶ月", 10, MID_GRAY, False),
        ("", 6, MID_GRAY, False),
        ("「最初は不安でしたが、スタッフさんが本当に優しくて安心できました。個室待機なので気まずさもなく、自分のペースで働けています！」", 10, DARK_GRAY, False),
    ], Cm(4.2), y1, Cm(6.0), Cm(3.5))

    add_divider(slide, y1 + Cm(4.0))

    # インタビュー2
    y2 = Cm(7.3)
    add_image_placeholder(slide, Cm(6.8), y2, Cm(3.0), Cm(3.5),
                          "女性写真2\n（モザイク）")

    add_multiline_text(slide, [
        ("Dさん（30歳・主婦）", 12, PURPLE, True),
        ("入店6ヶ月", 10, MID_GRAY, False),
        ("", 6, MID_GRAY, False),
        ("「子どもの送り迎えの間に働いています。自由出勤なので予定に合わせやすいし、身バレ対策もしっかりしていて安心です。」", 10, DARK_GRAY, False),
    ], Cm(0.8), y2, Cm(5.6), Cm(3.5))

    add_divider(slide, y2 + Cm(4.0))

    # インタビュー3
    y3 = Cm(11.8)
    add_image_placeholder(slide, Cm(0.8), y3, Cm(3.0), Cm(3.5),
                          "女性写真3\n（モザイク）")

    add_multiline_text(slide, [
        ("Eさん（26歳・移籍組）", 12, PURPLE, True),
        ("入店1年", 10, MID_GRAY, False),
        ("", 6, MID_GRAY, False),
        ("「前のお店では稼げなかったけど、ここはお客さんが多い！大手グループだから集客力が全然違います。待遇面もしっかりしています。」", 10, DARK_GRAY, False),
    ], Cm(4.2), y3, Cm(6.0), Cm(3.5))

    # 口コミ件数
    add_text_box(slide, "口コミ 99+ 件掲載中！",
                 Cm(0.8), Cm(16.0), Cm(9.2), Cm(0.7),
                 font_size=12, font_color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    # CTA
    add_rounded_rect(slide, Cm(1.5), Cm(17.0), Cm(7.8), Cm(1.2),
                     PURPLE, "もっと口コミを見る ▶", font_size=13, bold=True)

    return slide


# ============================================================
# スライド7: 未経験者向けセクション
# ============================================================
def slide_beginners(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, WHITE)

    # ヘッダー
    header_bg = add_bg_rect(slide, PINK_MAIN, 0, 0, SLIDE_WIDTH, Cm(2.8))
    add_multiline_text(slide, [
        ("未経験さん限定！", 20, WHITE, True),
        ("特別特典あり", 16, RGBColor(0xFF, 0xD5, 0x4F), True),
    ], Cm(0.8), Cm(0.3), Cm(9.2), Cm(2.2), alignment=PP_ALIGN.CENTER)

    # 体験保証
    add_rounded_rect(slide, Cm(1.5), Cm(3.4), Cm(7.8), Cm(1.5),
                     GOLD, "4時間 40,000円 体験保証", font_size=16, bold=True)

    # キービジュアル
    add_image_placeholder(slide, Cm(0.8), Cm(5.4), Cm(9.2), Cm(3.0),
                          "未経験者向けイメージ画像")

    # オススメポイント
    y = Cm(8.8)
    add_text_box(slide, "未経験さんへの オススメポイント",
                 Cm(0.8), y, Cm(9.2), Cm(0.8),
                 font_size=14, font_color=PINK_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    points = [
        "マンツーマン面接で何でも相談OK",
        "講習制度あり（接客マナーから丁寧に）",
        "最初から高収入が期待できる",
        "先輩スタッフが親身にサポート",
        "合わなければ即日退店OK",
        "ノルマ・罰金一切なし",
        "顔出し不要で身バレの心配なし",
    ]

    for i, point in enumerate(points):
        py = y + Cm(1.2) + Cm(i * 1.1)
        add_text_box(slide, f"✦ {point}", Cm(1.2), py, Cm(8.5), Cm(0.8),
                     font_size=11, font_color=DARK_GRAY, bold=False)

    # CTA
    add_rounded_rect(slide, Cm(1.5), Cm(17.5), Cm(7.8), Cm(1.2),
                     PINK_MAIN, "未経験でも安心！応募する ▶", font_size=13, bold=True)

    return slide


# ============================================================
# スライド8: 店内環境・設備
# ============================================================
def slide_facility(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, LIGHT_GRAY)

    # ヘッダー
    add_text_box(slide, "お店の環境", Cm(0.8), Cm(0.8), Cm(9.2), Cm(1.0),
                 font_size=20, font_color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, "清潔で快適な空間をご用意しています",
                 Cm(0.8), Cm(1.8), Cm(9.2), Cm(0.6),
                 font_size=11, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # 写真ギャラリー（2x3グリッド）
    photo_labels = [
        "店内写真", "待機室", "シャワールーム",
        "ロッカー", "受付", "スタッフ写真"
    ]

    for i, label in enumerate(photo_labels):
        col = i % 2
        row = i // 2
        x = Cm(0.8) + Cm(col * 4.8)
        y = Cm(2.8) + Cm(row * 3.5)
        add_image_placeholder(slide, x, y, Cm(4.4), Cm(3.0), label)

    # 設備リスト
    y_list = Cm(13.5)
    add_text_box(slide, "充実の設備", Cm(0.8), y_list, Cm(9.2), Cm(0.7),
                 font_size=13, font_color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    facilities = [
        "Wi-Fi完備 ／ ヘアアイロン・ドライヤー完備",
        "シャワールーム完備 ／ 鍵付きロッカー（1人1個）",
        "殺菌ソープ・うがい薬完備",
        "個室待機（テレビ・漫画あり）",
    ]
    for i, f in enumerate(facilities):
        add_text_box(slide, f"● {f}", Cm(0.8), y_list + Cm(0.9) + Cm(i * 0.9), Cm(9.2), Cm(0.7),
                     font_size=10, font_color=MID_GRAY)

    return slide


# ============================================================
# スライド9: 1日のスケジュール例
# ============================================================
def slide_schedule(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, WHITE)

    # ヘッダー
    header_bg = add_bg_rect(slide, PINK_MAIN, 0, 0, SLIDE_WIDTH, Cm(2.2))
    add_text_box(slide, "1日のスケジュール例", Cm(0.8), Cm(0.5), Cm(9.2), Cm(1.2),
                 font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    schedule = [
        ("10:00", "出勤・準備", "お店に到着したら個室待機室へ。身支度を整えます。"),
        ("10:30", "1本目のお客様", "60分コース。スタッフが近くに常駐しているので安心。"),
        ("12:00", "休憩", "自由時間。スマホや漫画で過ごせます。"),
        ("13:00", "2本目のお客様", "指名のリピーター様。"),
        ("14:30", "3本目のお客様", "新規のお客様。"),
        ("16:00", "4本目のお客様", "最後のお客様。"),
        ("17:00", "退勤", "お疲れ様でした！日払いで即日精算も可能。"),
    ]

    y = Cm(2.8)
    for i, (time_str, title, desc) in enumerate(schedule):
        row_y = y + Cm(i * 2.2)

        # 時間ラベル
        add_rounded_rect(slide, Cm(0.5), row_y, Cm(2.0), Cm(0.8),
                         PINK_MAIN, time_str, font_size=10, bold=True)

        # タイトル
        add_text_box(slide, title, Cm(2.8), row_y, Cm(7.0), Cm(0.7),
                     font_size=12, font_color=DARK_GRAY, bold=True)

        # 説明
        add_text_box(slide, desc, Cm(2.8), row_y + Cm(0.8), Cm(7.0), Cm(1.0),
                     font_size=9, font_color=MID_GRAY)

    # 合計収入
    add_rounded_rect(slide, Cm(1.0), Cm(18.0), Cm(8.8), Cm(0.8),
                     GOLD, "この日の収入：約51,000円（4本接客）", font_size=12, bold=True)

    return slide


# ============================================================
# スライド10: 面接の流れ
# ============================================================
def slide_flow(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, LIGHT_GRAY)

    add_text_box(slide, "面接〜体験入店の流れ", Cm(0.8), Cm(0.8), Cm(9.2), Cm(1.0),
                 font_size=18, font_color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    steps = [
        ("STEP 1", "応募", "LINE・電話・メールで\nお気軽にご連絡ください", PINK_MAIN),
        ("STEP 2", "面接", "マンツーマンで丁寧にご説明\n質問は何でもOK！", RGBColor(0xE9, 0x5E, 0x63)),
        ("STEP 3", "体験入店", "実際に体験してお店の雰囲気を\n確認してください", PURPLE),
        ("STEP 4", "本入店", "気に入っていただけたら\n本入店へ！即日入店もOK", GOLD),
    ]

    for i, (step_label, title, desc, color) in enumerate(steps):
        y = Cm(2.5) + Cm(i * 3.8)

        # ステップラベル
        add_rounded_rect(slide, Cm(1.5), y, Cm(3.0), Cm(0.8),
                         color, step_label, font_size=11, bold=True)

        # タイトル
        add_text_box(slide, title, Cm(5.0), y, Cm(4.8), Cm(0.8),
                     font_size=16, font_color=color, bold=True)

        # 説明
        add_text_box(slide, desc, Cm(1.5), y + Cm(1.0), Cm(7.8), Cm(1.5),
                     font_size=11, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

        # 矢印（最後以外）
        if i < len(steps) - 1:
            arrow_y = y + Cm(2.8)
            add_text_box(slide, "▼", Cm(0.8), arrow_y, Cm(9.2), Cm(0.8),
                         font_size=16, font_color=RGBColor(0xCC, 0xCC, 0xCC),
                         alignment=PP_ALIGN.CENTER)

    # 補足
    add_text_box(slide, "※面接時の交通費支給！\n※見学のみでもOKです！",
                 Cm(0.8), Cm(17.5), Cm(9.2), Cm(1.2),
                 font_size=11, font_color=PINK_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================
# スライド11: アクセス・店舗情報
# ============================================================
def slide_access(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, WHITE)

    # ヘッダー
    header_bg = add_bg_rect(slide, DARK_GRAY, 0, 0, SLIDE_WIDTH, Cm(2.2))
    add_text_box(slide, "店舗情報・アクセス", Cm(0.8), Cm(0.5), Cm(9.2), Cm(1.2),
                 font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 地図プレースホルダー
    add_image_placeholder(slide, Cm(0.8), Cm(2.8), Cm(9.2), Cm(4.5),
                          "Google Map / アクセスマップ画像")

    # 店舗情報
    info = [
        ("店名", "JJクラブ 大曽根"),
        ("業種", "デリヘル"),
        ("エリア", "愛知県 名古屋市 大曽根"),
        ("最寄り駅", "JR大曽根駅・地下鉄大曽根駅"),
        ("営業時間", "10:00〜24:00"),
        ("定休日", "年中無休"),
        ("勤務時間", "最短3時間〜OK"),
        ("出勤", "完全自由出勤"),
        ("応募資格", "18歳〜30代前半（高校生NG）"),
    ]

    y = Cm(7.8)
    for i, (label, value) in enumerate(info):
        row_y = y + Cm(i * 1.1)

        # ラベル
        add_rounded_rect(slide, Cm(0.8), row_y, Cm(2.8), Cm(0.8),
                         PINK_LIGHT, label, font_size=9,
                         font_color=PINK_DARK, bold=True)

        # 値
        add_text_box(slide, value, Cm(4.0), row_y, Cm(6.0), Cm(0.8),
                     font_size=11, font_color=DARK_GRAY, bold=False)

    return slide


# ============================================================
# スライド12: こだわり条件タグ
# ============================================================
def slide_tags(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, PINK_LIGHT)

    add_text_box(slide, "こだわり条件", Cm(0.8), Cm(0.8), Cm(9.2), Cm(1.0),
                 font_size=18, font_color=PINK_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, "JJクラブ大曽根はこんなお店です",
                 Cm(0.8), Cm(1.8), Cm(9.2), Cm(0.6),
                 font_size=11, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    tags = [
        "個室待機OK", "自宅待機OK", "送迎あり", "日払い",
        "アリバイ対策", "顔出しなし", "体験入店あり", "出張面接",
        "交通費支給", "制服貸与", "ノルマなし", "罰金なし",
        "未経験歓迎", "ブランクOK", "掛け持ちOK", "短期OK",
        "友達と応募OK", "Wi-Fi完備", "寮完備", "託児所紹介",
        "ヘアメイク", "写真撮影", "衛生対策", "スタッフ常駐",
    ]

    y = Cm(2.8)
    x = Cm(0.5)
    tag_w = Cm(3.0)
    tag_h = Cm(0.9)
    gap_x = Cm(0.3)
    gap_y = Cm(0.3)
    cols = 3

    for i, tag in enumerate(tags):
        col = i % cols
        row = i // cols
        tx = x + Cm(col * 3.3)
        ty = y + Cm(row * 1.2)

        add_rounded_rect(slide, tx, ty, tag_w, tag_h,
                         WHITE, tag, font_size=9, font_color=PINK_DARK, bold=False)

    # キービジュアル
    add_image_placeholder(slide, Cm(0.8), Cm(12.5), Cm(9.2), Cm(3.5),
                          "お店の魅力が伝わる\nイメージ画像")

    # CTA
    add_rounded_rect(slide, Cm(1.5), Cm(16.5), Cm(7.8), Cm(1.2),
                     PINK_MAIN, "詳しく聞いてみる ▶", font_size=13, bold=True)

    return slide


# ============================================================
# スライド13: 応募方法（CTA集約）
# ============================================================
def slide_apply(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, BLACK)

    # ヘッダー
    add_multiline_text(slide, [
        ("今すぐ応募！", 24, WHITE, True),
        ("お気軽にお問い合わせください", 13, RGBColor(0xBB, 0xBB, 0xBB), False),
    ], Cm(0.8), Cm(1.0), Cm(9.2), Cm(2.5), alignment=PP_ALIGN.CENTER)

    # LINE応募
    y = Cm(3.5)
    add_rounded_rect(slide, Cm(1.0), y, Cm(8.8), Cm(2.5),
                     GREEN_LINE, "", font_size=14, bold=True)
    add_multiline_text(slide, [
        ("LINE で応募", 16, WHITE, True),
        ("ID: showtikubai", 12, WHITE, False),
        ("「求人を見ました」とメッセージ♪", 10, RGBColor(0xCC, 0xFF, 0xCC), False),
    ], Cm(1.5), y + Cm(0.2), Cm(5.5), Cm(2.2))
    add_image_placeholder(slide, Cm(7.5), y + Cm(0.3), Cm(1.8), Cm(1.8), "QRコード")

    # 電話応募
    y2 = Cm(6.5)
    add_rounded_rect(slide, Cm(1.0), y2, Cm(8.8), Cm(2.0),
                     PINK_MAIN, "", font_size=14, bold=True)
    add_multiline_text(slide, [
        ("電話で応募", 16, WHITE, True),
        ("24時間受付中！", 11, RGBColor(0xFF, 0xCC, 0xDD), False),
        ("「求人を見ました」とお伝えください", 10, RGBColor(0xFF, 0xCC, 0xDD), False),
    ], Cm(1.5), y2 + Cm(0.2), Cm(7.5), Cm(1.8))

    # メール応募
    y3 = Cm(9.0)
    add_rounded_rect(slide, Cm(1.0), y3, Cm(8.8), Cm(1.5),
                     PURPLE, "メールで応募する ✉", font_size=14, bold=True)

    # 応募特典
    y4 = Cm(11.0)
    bonus_bg = add_bg_rect(slide, GOLD, Cm(1.0), y4, Cm(8.8), Cm(3.0))
    add_multiline_text(slide, [
        ("応募特典！", 16, WHITE, True),
        ("", 6, WHITE, False),
        ("面接時の交通費支給！", 14, WHITE, True),
        ("「求人を見ました」と伝えるだけ♪", 11, WHITE, False),
    ], Cm(1.5), y4 + Cm(0.2), Cm(7.8), Cm(2.8), alignment=PP_ALIGN.CENTER)

    # 求人内容に嘘偽りなし
    add_multiline_text(slide, [
        ("求人内容に嘘や偽りは一切ありません。", 11, RGBColor(0xBB, 0xBB, 0xBB), False),
        ("見学・面接・体験等で", 11, RGBColor(0xBB, 0xBB, 0xBB), False),
        ("お店を知ってもらい、", 11, RGBColor(0xBB, 0xBB, 0xBB), False),
        ("ご納得の上でご入店ください。", 11, RGBColor(0xBB, 0xBB, 0xBB), False),
    ], Cm(0.8), Cm(14.5), Cm(9.2), Cm(3.0), alignment=PP_ALIGN.CENTER)

    # グループ紹介
    add_text_box(slide, "JJグループ（新栄・池下・大曽根・堀田）",
                 Cm(0.8), Cm(17.0), Cm(9.2), Cm(0.7),
                 font_size=10, font_color=RGBColor(0x99, 0x99, 0x99),
                 alignment=PP_ALIGN.CENTER)

    # キービジュアル
    add_image_placeholder(slide, Cm(1.5), Cm(17.8), Cm(7.8), Cm(1.0),
                          "グループロゴ / バナー")

    return slide


# ============================================================
# スライド14: 応募フォーム案内
# ============================================================
def slide_form(prs, layout):
    slide = prs.slides.add_slide(layout)
    add_bg_rect(slide, WHITE)

    # ヘッダー
    header_bg = add_bg_rect(slide, PINK_MAIN, 0, 0, SLIDE_WIDTH, Cm(2.2))
    add_text_box(slide, "カンタン応募フォーム", Cm(0.8), Cm(0.5), Cm(9.2), Cm(1.2),
                 font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # フォーム項目
    fields = [
        "お名前（ニックネームOK）",
        "年齢",
        "お住まいの地域",
        "風俗経験",
        "ご希望の面接日時",
        "働き方の希望",
        "メッセージ（任意）",
    ]

    y = Cm(3.0)
    for i, field in enumerate(fields):
        fy = y + Cm(i * 1.8)

        # ラベル
        add_text_box(slide, field, Cm(1.0), fy, Cm(8.8), Cm(0.6),
                     font_size=11, font_color=DARK_GRAY, bold=True)

        # 入力欄風の矩形
        input_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.0), fy + Cm(0.6), Cm(8.8), Cm(0.9)
        )
        input_box.fill.solid()
        input_box.fill.fore_color.rgb = LIGHT_GRAY
        input_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        input_box.line.width = Pt(1)

    # 送信ボタン
    add_rounded_rect(slide, Cm(1.5), Cm(16.5), Cm(7.8), Cm(1.3),
                     PINK_MAIN, "この内容で応募する", font_size=16, bold=True)

    add_text_box(slide, "※ 18歳未満の方はご応募いただけません",
                 Cm(0.8), Cm(18.0), Cm(9.2), Cm(0.6),
                 font_size=9, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================
# メイン実行
# ============================================================
def main():
    prs, layout = create_presentation()

    # 全スライドを生成
    slide_hero(prs, layout)        # 1. ファーストビュー
    slide_merits(prs, layout)      # 2. 3つのメリット
    slide_trial(prs, layout)       # 3. 体験入店情報
    slide_safety(prs, layout)      # 4. 安心ポイント
    slide_salary(prs, layout)      # 5. 給与情報
    slide_interview(prs, layout)   # 6. 在籍女性の声
    slide_beginners(prs, layout)   # 7. 未経験者向け
    slide_facility(prs, layout)    # 8. 店内環境
    slide_schedule(prs, layout)    # 9. 1日のスケジュール
    slide_flow(prs, layout)        # 10. 面接の流れ
    slide_access(prs, layout)      # 11. アクセス・店舗情報
    slide_tags(prs, layout)        # 12. こだわり条件
    slide_apply(prs, layout)       # 13. 応募方法
    slide_form(prs, layout)        # 14. 応募フォーム

    # 出力
    output_path = os.path.join(os.path.dirname(__file__), 'JJ大曽根_LP_デザイン.pptx')
    prs.save(output_path)
    print(f"PPTX saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
